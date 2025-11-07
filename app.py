import os
import json
import telebot
import requests
from telebot import types
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
from PIL import Image
import logging

# Logging sozlash
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# 🔐 Tokenlar
BOT_TOKEN = "8142445503:AAHoF0S0of-bCQmPkOitpHHg9KXTL51bPIg"
GEMINI_API_KEY = "AIzaSyAs3WE0NyNn_uuNQG2KVUA_deQbYGEOG-8"
UNSPLASH_ACCESS_KEY = "yK3fWW35A1cS0h3yV0aq0PNVNyaK62e4V9cENRnDbaQ"

bot = telebot.TeleBot(BOT_TOKEN)


# ========== RASM QIDIRISH (UNSPLASH) ==========
def search_image(keyword):
    """Unsplash API orqali mavzuga mos rasm topish"""
    if not UNSPLASH_ACCESS_KEY or UNSPLASH_ACCESS_KEY == "UNSPLASH_KEY_SINI_BU_YERGA":
        logger.warning("Unsplash API key yo'q")
        return None
    
    try:
        url = "https://api.unsplash.com/search/photos"
        params = {
            "query": keyword,
            "per_page": 1,
            "orientation": "landscape",
            "client_id": UNSPLASH_ACCESS_KEY
        }
        
        response = requests.get(url, params=params, timeout=10)
        response.raise_for_status()
        data = response.json()
        
        if data.get("results") and len(data["results"]) > 0:
            image_url = data["results"][0]["urls"]["regular"]
            logger.info(f"Rasm topildi: {keyword}")
            return image_url
        else:
            logger.warning(f"Rasm topilmadi: {keyword}")
            return None
        
    except requests.exceptions.RequestException as e:
        logger.error(f"Rasm qidirish xatosi: {e}")
        return None
    except Exception as e:
        logger.error(f"Kutilmagan xato: {e}")
        return None


def download_image(url):
    """Rasmni yuklab olish va BytesIO ga saqlash"""
    try:
        response = requests.get(url, timeout=15)
        response.raise_for_status()
        
        img = Image.open(BytesIO(response.content))
        
        # RGB ga o'tkazish
        if img.mode in ('RGBA', 'LA', 'P'):
            background = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            if img.mode in ('RGBA', 'LA'):
                background.paste(img, mask=img.split()[-1])
            else:
                background.paste(img)
            img = background
        
        # O'lchamni kichraytirish
        max_size = (1600, 900)
        img.thumbnail(max_size, Image.Resampling.LANCZOS)
        
        # BytesIO ga saqlash
        img_byte_arr = BytesIO()
        img.save(img_byte_arr, format='JPEG', quality=85, optimize=True)
        img_byte_arr.seek(0)
        
        logger.info("Rasm muvaffaqiyatli yuklandi")
        return img_byte_arr
        
    except Exception as e:
        logger.error(f"Rasm yuklash xatosi: {e}")
        return None


# ========== GEMINI AI FUNKSIYASI ==========
def generate_ppt_content(topic):
    """Gemini AI yordamida JSON formatda slaydlar kontenti yaratish"""
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash-exp:generateContent?key={GEMINI_API_KEY}"

    prompt = f"""
    Sen professional taqdimot dizaynerisan.
    Mavzu: "{topic}"
    Jami 6-8 ta slayd yarating.
    
    MUHIM: Faqat JSON formatda javob ber, boshqa hech narsa yozma!
    
    Quyidagi JSON formatini qat'iy rioya qil:
    [
      {{
        "sarlavha": "Kirish",
        "matn": "Mavzuga kirish matni. 3-4 ta qisqa va tushunarli jumla.",
        "rasm_kalit": "introduction concept business"
      }},
      {{
        "sarlavha": "Asosiy Nuqtalar",
        "matn": "Muhim faktlar va ma'lumotlar. 4-5 ta jumla.",
        "rasm_kalit": "key points strategy"
      }},
      {{
        "sarlavha": "Xulosa",
        "matn": "Yakuniy fikrlar va takliflar. 3-4 ta jumla.",
        "rasm_kalit": "conclusion success future"
      }}
    ]
    
    QOIDALAR:
    - Har bir sarlavha 2-5 so'zdan iborat
    - Matn aniq, qisqa, foydali (har birida 3-5 jumla)
    - rasm_kalit - ingliz tilida 2-3 ta kalit so'z
    - Faqat JSON, boshqa hech narsa!
    """

    headers = {"Content-Type": "application/json"}
    payload = {
        "contents": [{
            "parts": [{"text": prompt}]
        }],
        "generationConfig": {
            "temperature": 0.7,
            "topK": 40,
            "topP": 0.95,
            "maxOutputTokens": 2048,
        }
    }

    try:
        response = requests.post(url, headers=headers, json=payload, timeout=40)
        response.raise_for_status()
        data = response.json()
        
        if "candidates" not in data or len(data["candidates"]) == 0:
            logger.error("Gemini javob bermadi")
            return None
        
        text = data["candidates"][0]["content"]["parts"][0]["text"]
        text = text.strip()
        
        # JSON tozalash
        text = text.replace("```json", "").replace("```", "").strip()
        
        # JSON parse
        slides = json.loads(text)
        
        if not isinstance(slides, list) or len(slides) < 3:
            logger.error(f"Slaydlar soni kam: {len(slides) if isinstance(slides, list) else 0}")
            return None
        
        logger.info(f"AI kontent tayyor: {len(slides)} ta slayd")
        return slides
        
    except json.JSONDecodeError as e:
        logger.error(f"JSON parse xatosi: {e}")
        logger.error(f"Matn: {text[:200] if 'text' in locals() else 'yo`q'}")
        return None
    except requests.exceptions.RequestException as e:
        logger.error(f"API so'rov xatosi: {e}")
        return None
    except Exception as e:
        logger.error(f"AI kontent xatosi: {e}")
        return None


# ========== POWERPOINT YARATISH ==========
def create_ppt(topic, slides, with_images=True):
    """AI kontent asosida rasmli PowerPoint yaratish"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        for i, s in enumerate(slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Fon rangi
            background = slide.background
            fill = background.fill
            fill.solid()
            
            colors = [
                RGBColor(245, 247, 255),
                RGBColor(255, 248, 240),
                RGBColor(240, 255, 245),
            ]
            fill.fore_color.rgb = colors[i % 3]

            # RASM QO'SHISH
            image_added = False
            if with_images and s.get("rasm_kalit"):
                try:
                    logger.info(f"Rasm qidirilmoqda: {s.get('rasm_kalit')}")
                    image_url = search_image(s.get("rasm_kalit"))
                    
                    if image_url:
                        img_stream = download_image(image_url)
                        if img_stream:
                            left = Inches(5.5)
                            top = Inches(2)
                            width = Inches(4)
                            height = Inches(3)
                            
                            slide.shapes.add_picture(
                                img_stream, left, top,
                                width=width, height=height
                            )
                            image_added = True
                            logger.info(f"Rasm qo'shildi: slayd {i+1}")
                except Exception as e:
                    logger.warning(f"Rasm qo'shishda xato (slayd {i+1}): {e}")

            # Yuqori chiziq
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(0), Inches(0), 
                Inches(10), Inches(0.25)
            )
            fill_shape = shape.fill
            fill_shape.solid()
            fill_shape.fore_color.rgb = RGBColor(0, 120, 212)
            shape.line.fill.background()

            # Sarlavha
            title_width = Inches(4.5) if image_added else Inches(9)
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.6),
                title_width, Inches(1)
            )
            tf = title_box.text_frame
            tf.text = s.get("sarlavha", "").strip()
            tf.word_wrap = True

            if tf.paragraphs and tf.paragraphs[0].runs:
                title_run = tf.paragraphs[0].runs[0]
                title_run.font.bold = True
                title_run.font.size = Pt(36)
                title_run.font.color.rgb = RGBColor(25, 25, 112)

            # Matn
            body_width = Inches(4.5) if image_added else Inches(9)
            body_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8),
                body_width, Inches(5)
            )
            tf_body = body_box.text_frame
            tf_body.text = s.get("matn", "").strip()
            tf_body.word_wrap = True
            
            for paragraph in tf_body.paragraphs:
                paragraph.space_before = Pt(8)
                paragraph.space_after = Pt(8)
                paragraph.line_spacing = 1.4
                
                for run in paragraph.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(18)
                    run.font.color.rgb = RGBColor(40, 40, 40)

            # Slayd raqami
            page_num_box = slide.shapes.add_textbox(
                Inches(9), Inches(7), 
                Inches(0.8), Inches(0.4)
            )
            tf_num = page_num_box.text_frame
            tf_num.text = f"{i + 1}/{len(slides)}"
            
            if tf_num.paragraphs and tf_num.paragraphs[0].runs:
                for run in tf_num.paragraphs[0].runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = RGBColor(120, 120, 120)

        # Fayl saqlash
        safe_filename = "".join(c for c in topic if c.isalnum() or c in (' ', '-', '_')).strip()
        safe_filename = safe_filename[:50] or "presentation"
        filename = f"{safe_filename}.pptx"
        
        prs.save(filename)
        logger.info(f"PPT yaratildi: {filename}")
        return filename
        
    except Exception as e:
        logger.error(f"PPT yaratishda xato: {e}")
        raise


# ========== TELEGRAM BOT ==========
@bot.message_handler(commands=['start', 'help'])
def start(message):
    try:
        first_name = message.from_user.first_name or ""
        last_name = message.from_user.last_name or ""
        user_name = f"{first_name} {last_name}".strip() or "Foydalanuvchi"
        
        welcome_text = (
            f"👋 Salom, {user_name}!\n\n"
            "Men — sun'iy intellekt asosidagi slayd tayyorlovchi botman. 🎨⚡\n\n"
            "💡 Nimalarga qodirman:\n"
            "• Mavzuni tahlil qilib, chiroyli matn yarataman\n"
            "• Har bir slaydga professional rasm qo'shaman 🖼️\n"
            "• 6-8 ta slayd tayyorlayman\n"
            "• Tayyor PowerPoint faylni yuboraman\n\n"
            "Boshlaymizmi? 👇"
        )

        markup = types.ReplyKeyboardMarkup(resize_keyboard=True, row_width=2)
        markup.add("📑 Rasmli slayd", "📄 Rasmsiz slayd")
        markup.add("❓ Yordam")

        bot.send_message(message.chat.id, welcome_text, reply_markup=markup)
        logger.info(f"Start buyrug'i: {message.from_user.id}")
        
    except Exception as e:
        logger.error(f"Start xatosi: {e}")
        bot.send_message(message.chat.id, "⚠️ Xatolik yuz berdi. Qaytadan /start bosing.")


@bot.message_handler(func=lambda msg: msg.text == "❓ Yordam")
def help_message(message):
    help_text = (
        "📖 *Qo'llanma:*\n\n"
        "1️⃣ Slayd turini tanlang:\n"
        "   • 📑 *Rasmli slayd* - har bir slaydda rasm\n"
        "   • 📄 *Rasmsiz slayd* - faqat matn\n\n"
        "2️⃣ Mavzu nomini yozing\n"
        "3️⃣ 10-20 soniya kuting\n"
        "4️⃣ Tayyor faylni yuklab oling!\n\n"
        "💡 *Maslahatlar:*\n"
        "• Aniq mavzu yozing\n"
        "• Rasmli versiya biroz ko'proq vaqt oladi\n"
        "• O'zbek yoki ingliz tilida ishlaydi\n\n"
        "❓ Muammo bo'lsa /start bosing yoki @vafoyev_n ga yozing"
    )
    bot.send_message(message.chat.id, help_text, parse_mode="Markdown")


@bot.message_handler(func=lambda msg: msg.text in ["📑 Rasmli slayd", "📄 Rasmsiz slayd"])
def ask_topic(message):
    with_images = (message.text == "📑 Rasmli slayd")
    
    msg = bot.send_message(
        message.chat.id,
        f"{'🖼️ Rasmli' if with_images else '📄 Rasmsiz'} slayd yaratamiz!\n\n"
        "🧠 Qaysi mavzu bo'yicha?\n"
        "Mavzuni yozing ✍️"
    )
    bot.register_next_step_handler(msg, process_topic, with_images)


def process_topic(message, with_images=True):
    try:
        topic = message.text.strip()
        
        # Validatsiya
        if len(topic) < 3:
            bot.send_message(message.chat.id, "⚠️ Mavzu juda qisqa. Kamida 3 ta belgi yozing.")
            return
        
        if len(topic) > 200:
            bot.send_message(message.chat.id, "⚠️ Mavzu juda uzun. Maksimal 200 belgi.")
            return
        
        logger.info(f"Mavzu: '{topic}', Foydalanuvchi: {message.from_user.id}, Rasmli: {with_images}")
        
        processing_msg = bot.send_message(
            message.chat.id,
            f"⚡ *{topic}* uchun taqdimot tayyorlanmoqda...\n\n"
            "🤖 AI kontent yaratmoqda...",
            parse_mode="Markdown"
        )

        # AI kontent
        slides = generate_ppt_content(topic)
        if not slides:
            bot.edit_message_text(
                "⚠️ AI kontent yaratishda xatolik yuz berdi.\n"
                "Iltimos, mavzuni boshqacha yozib qaytadan urinib ko'ring.",
                message.chat.id, processing_msg.message_id
            )
            return

        bot.edit_message_text(
            f"✅ Kontent tayyor! ({len(slides)} ta slayd)\n\n"
            f"{'🖼️ Rasmlar yuklanmoqda...' if with_images else '🎨 Dizayn qo`llanmoqda...'}",
            message.chat.id, processing_msg.message_id
        )
        
        # PPT yaratish
        ppt_file = create_ppt(topic, slides, with_images)
        
        if not ppt_file or not os.path.exists(ppt_file):
            bot.edit_message_text(
                "⚠️ Fayl yaratishda xatolik. Qaytadan urinib ko'ring.",
                message.chat.id, processing_msg.message_id
            )
            return
        
        bot.edit_message_text(
            "📤 Fayl yuborilmoqda...",
            message.chat.id, processing_msg.message_id
        )
        
        # Faylni yuborish
        with open(ppt_file, "rb") as f:
            bot.send_document(
                message.chat.id, f,
                caption=f"✅ *{topic}* mavzusidagi taqdimot tayyor!\n\n"
                        f"📊 Slaydlar: {len(slides)} ta\n"
                        f"{'🖼️ Rasmlar qo`shildi' if with_images else '📄 Rasmsiz versiya'}",
                parse_mode="Markdown"
            )

        # Tozalash
        try:
            os.remove(ppt_file)
            bot.delete_message(message.chat.id, processing_msg.message_id)
        except:
            pass
        
        # Keyingi harakat
        markup = types.ReplyKeyboardMarkup(resize_keyboard=True, row_width=2)
        markup.add("📑 Rasmli slayd", "📄 Rasmsiz slayd")
        markup.add("❓ Yordam")
        bot.send_message(message.chat.id, "✨ Yana taqdimot yarataymizmi?", reply_markup=markup)
        
        logger.info(f"Taqdimot yuborildi: {topic}")
        
    except Exception as e:
        logger.error(f"Process xatosi: {e}", exc_info=True)
        bot.send_message(
            message.chat.id,
            "⚠️ Xatolik yuz berdi. Iltimos, qaytadan urinib ko'ring.\n"
            "Agar muammo takrorlansa, @vafoyev_n ga xabar bering."
        )


@bot.message_handler(func=lambda msg: True)
def handle_text(message):
    bot.send_message(
        message.chat.id,
        "🤔 Tushunmadim. Tugmalardan birini tanlang yoki /start bosing."
    )


# ========== ISHGA TUSHIRISH ==========
if __name__ == "__main__":
    print("=" * 50)
    print("🚀 Bot ishga tushdi!")
    print("🖼️ Rasmli slaydlar rejimi: FAOL")
    print("📍 To'xtatish: Ctrl+C")
    print("=" * 50)
    
    # API keys tekshirish
    if not BOT_TOKEN or BOT_TOKEN == "YOUR_BOT_TOKEN":
        print("❌ XATO: BOT_TOKEN kiritilmagan!")
        exit(1)
    
    if not GEMINI_API_KEY:
        print("❌ XATO: GEMINI_API_KEY kiritilmagan!")
        exit(1)
    
    if not UNSPLASH_ACCESS_KEY or UNSPLASH_ACCESS_KEY == "UNSPLASH_KEY_SINI_BU_YERGA":
        print("⚠️ OGOHLANTIRISH: Unsplash API key yo'q")
        print("📝 https://unsplash.com/developers dan oling")
        print("🔧 Hozircha rasmsiz rejimda ishlaydi\n")
    
    try:
        logger.info("Bot ishga tushirilmoqda...")
        bot.infinity_polling(timeout=60, long_polling_timeout=60)
    except KeyboardInterrupt:
        print("\n👋 Bot to'xtatildi")
        logger.info("Bot to'xtatildi")
    except Exception as e:
        print(f"❌ Bot xatosi: {e}")
        logger.error(f"Bot xatosi: {e}", exc_info=True)